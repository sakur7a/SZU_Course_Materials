from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "outputs" / "map_coloring_oral_slides.pptx"

BLUE = RGBColor(31, 78, 121)
GRAY = RGBColor(75, 85, 99)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)

FONT = "Microsoft YaHei"
FONT_LATIN = "Aptos"


def set_run(run, size=20, color=BLACK, bold=False, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_text(slide, text, x, y, w, h, size=20, color=BLACK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    set_run(r, size=size, color=color, bold=bold)
    return box


def add_multiline(slide, lines, x, y, w, h, size=18, color=BLACK, bullet=False, numbered=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        if bullet:
            p.level = 0
            p.text = "▶  " + line
        if numbered:
            p.text = f"{i + 1}.  {line}"
    return box


def add_title(slide, title):
    add_text(slide, title, 0.35, 0.25, 12.2, 0.45, size=25, color=BLUE)


def add_footer(slide, idx, total=7):
    add_text(slide, "Map Coloring Lab 3", 0.55, 7.15, 3.0, 0.25, size=9, color=GRAY)
    add_text(slide, f"{idx}/{total}", 12.25, 7.15, 0.55, 0.25, size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def set_shape_line(shape, color=BLUE, width=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    shape.fill.background()


def add_rule(slide, x, y, w, width=1.2, color=BLACK):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + w), Inches(y))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_three_line_table(slide, headers, rows, x, y, w, col_widths, row_h=0.35, font_size=14):
    add_rule(slide, x, y, w, 1.1)
    y_header = y + 0.07
    cursor = x
    for header, cw in zip(headers, col_widths):
        add_text(slide, header, cursor, y_header, cw, 0.25, size=font_size, bold=True)
        cursor += cw
    add_rule(slide, x, y + 0.38, w, 0.9)
    current_y = y + 0.45
    for row in rows:
        cursor = x
        for value, cw in zip(row, col_widths):
            add_text(slide, str(value), cursor, current_y, cw, row_h, size=font_size)
            cursor += cw
        current_y += row_h
    add_rule(slide, x, current_y + 0.03, w, 1.1)


def add_metric(slide, value, label, x):
    add_text(slide, value, x, 5.55, 2.2, 0.45, size=30, color=BLUE)
    add_text(slide, label, x, 6.0, 2.4, 0.35, size=14, color=GRAY)


def add_flow(slide):
    labels = ["贪心初始解", "定位冲突顶点", "选择最优非禁忌移动", "更新禁忌表", "冲突边数 = 0"]
    x, y, w, h = 8.0, 2.05, 3.4, 0.38
    centers = []
    for i, label in enumerate(labels):
        yy = y + i * 0.68
        rect = slide.shapes.add_shape(1, Inches(x), Inches(yy), Inches(w), Inches(h))
        set_shape_line(rect, BLUE, 1.0)
        tf = rect.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        set_run(r, size=12)
        centers.append((x + w / 2, yy + h))
    for i in range(len(centers) - 1):
        x1, y1 = centers[i]
        x2, y2 = centers[i + 1][0], y + (i + 1) * 0.68
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = BLUE
        line.line.width = Pt(1.0)


def add_cover_graph(slide):
    points = [(3.9, 4.8), (5.55, 4.35), (7.1, 4.75), (8.9, 4.2), (4.25, 5.42), (5.8, 5.32), (7.65, 5.48)]
    edges = [(0, 1), (1, 2), (2, 3), (4, 1), (1, 5), (5, 2), (2, 6), (6, 3)]
    for a, b in edges:
        x1, y1 = points[a]
        x2, y2 = points[b]
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = BLUE
        line.line.width = Pt(1.6)
    for i, (x, y) in enumerate(points):
        oval = slide.shapes.add_shape(9, Inches(x - 0.035), Inches(y - 0.035), Inches(0.07), Inches(0.07))
        oval.fill.solid()
        oval.fill.fore_color.rgb = BLUE if i < 4 else GRAY
        oval.line.color.rgb = BLUE if i < 4 else GRAY


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 cover
    s = prs.slides.add_slide(blank)
    add_text(s, "地图填色", 0.7, 1.05, 5.0, 0.75, size=39, color=BLUE, bold=True)
    add_text(s, "从暴力回溯到 DSATUR 与 TabuCol", 0.7, 1.9, 8.5, 0.55, size=27, bold=True)
    add_text(s, "一个可复现的图着色实验流水线：建模、剪枝优化、局部搜索、结果验证", 0.7, 3.0, 10.5, 0.45, size=18, color=GRAY)
    add_cover_graph(s)
    add_text(s, "白底极简版式；所有结果以冲突边数为合法性判据", 0.7, 6.55, 7.0, 0.3, size=13, color=GRAY)

    # 2 model
    s = prs.slides.add_slide(blank)
    add_title(s, "核心问题：把“地图”压缩成一个可验证模型")
    add_text(s, "不讲四色定理常识；直接讲我设计的实验对象与评价指标", 0.7, 1.05, 6.5, 0.3, size=14, color=GRAY)
    add_text(s, "区域邻接关系 → 图的 k-着色约束满足问题", 6.25, 1.0, 6.0, 0.55, size=25, color=BLUE, bold=True)
    add_text(s, "G = (V, E),    C = {1, …, k}\nf : V → C\n(u, v) ∈ E ⇒ f(u) ≠ f(v)", 1.8, 2.45, 4.0, 1.5, size=20)
    add_three_line_table(
        s,
        ["设计项", "本实验实现"],
        [["输入", "小图邻接、DIMACS .col、随机平面图"], ["目标", "给定颜色数下构造合法染色"], ["判据", "冲突边数 = 0"], ["输出", "每个顶点的颜色编号 CSV"]],
        6.6, 2.0, 5.1, [1.15, 3.95], row_h=0.36, font_size=15
    )
    add_metric(s, "0", "合法解冲突边数", 0.8)
    add_metric(s, "450", "基准图顶点数", 5.0)
    add_metric(s, "20 ~ 300", "平面图规模", 9.0)
    add_footer(s, 2)

    # 3 DSATUR
    s = prs.slides.add_slide(blank)
    add_title(s, "优化一：从固定顺序暴力回溯到 DSATUR")
    add_text(s, "复现 Brélaz 的 DSATUR 思想，并把它接入回溯与前向检查", 0.8, 1.05, 7.0, 0.3, size=15, color=GRAY)
    add_three_line_table(s, ["方法", "关键问题"], [["暴力回溯", "固定顶点顺序，最坏枚举 k^n"], ["DSATUR", "先处理约束最紧顶点"], ["前向检查", "某邻居无可用色时立即剪枝"]], 0.8, 1.7, 5.1, [1.8, 3.3], row_h=0.42, font_size=15)
    add_text(s, "饱和度定义", 0.85, 4.45, 2.2, 0.28, size=16, bold=True)
    add_text(s, "sat(v) = |{ f(u) | u ∈ N(v), f(u) ≠ ∅ }|", 1.0, 5.05, 5.5, 0.35, size=18)
    add_three_line_table(s, ["行", "DSATUR-Backtracking 核心流程"], [["1", "维护每个未着色顶点的可用颜色集合"], ["2", "选择 sat(v) 最大的未着色顶点"], ["3", "并列时选择度数 deg(v) 最大者"], ["4", "尝试颜色 c，并临时删除邻居域中的 c"], ["5", "若产生空颜色域，则立即回溯"], ["6", "若全部顶点着色完成，返回合法解"]], 6.35, 1.65, 5.5, [0.55, 4.95], row_h=0.33, font_size=13)
    add_text(s, "优化点不是改变问题，而是改变搜索顺序：更早遇到矛盾，更少进入无效分支。", 0.8, 6.22, 11.9, 0.52, size=22, color=BLUE, bold=True)
    add_footer(s, 3)

    # 4 TabuCol
    s = prs.slides.add_slide(blank)
    add_title(s, "优化二：450 顶点实例用 TabuCol 求可行解")
    add_text(s, "精确回溯不适合直接处理 le450；改用局部搜索快速构造 0 冲突解", 0.8, 1.1, 8.5, 0.3, size=15, color=GRAY)
    add_three_line_table(s, ["行", "TabuCol 局部搜索核心流程"], [["1", "贪心生成初始 k-着色，允许冲突"], ["2", "只在冲突顶点集合中选择移动"], ["3", "评价移动 Δ = 新冲突数 − 旧冲突数"], ["4", "最优移动若被禁忌，只有优于历史最优时允许"], ["5", "更新颜色与禁忌表，直到冲突边数为 0"]], 0.8, 1.75, 5.1, [0.55, 4.55], row_h=0.55, font_size=15)
    add_flow(s)
    add_text(s, "定位：DSATUR 做确定性验证；TabuCol 做大规模给定颜色数下的快速可行染色。", 0.8, 6.2, 10.7, 0.35, size=13, color=GRAY)
    add_footer(s, 4)

    # 5 comparison
    s = prs.slides.add_slide(blank)
    add_title(s, "复现实验：暴力回溯 vs DSATUR")
    add_text(s, "首解可能很快，但完整搜索空间暴露了暴力法的指数风险", 0.8, 1.25, 7.0, 0.3, size=15, color=GRAY)
    add_three_line_table(s, ["数据", "顶点", "暴力赋值", "DSATUR赋值", "暴力完整空间", "外推/s"], [["小规模地图", "9", "19", "9", "2.621×10^5", "2.580×10^-1"], ["平面图 n=20", "20", "51", "20", "1.100×10^12", "6.985×10^5"], ["平面图 n=32", "32", "78", "32", "1.845×10^19", "1.102×10^13"]], 0.8, 1.8, 10.6, [2.2, 0.8, 1.3, 1.45, 2.5, 2.35], row_h=0.34, font_size=13)
    add_text(s, "DSATUR 的优势体现在搜索\n节点数，而不只是一两个样例\n的运行时间。", 0.8, 4.1, 5.0, 1.4, size=24, color=BLUE, bold=True)
    add_multiline(s, ["暴力法：固定顺序，完整空间为 4^n", "DSATUR：动态选点，实验中赋值次数约等于 n", "外推只用于说明指数增长趋势，不当作实际运行时间"], 7.0, 4.05, 5.2, 1.75, size=15, bullet=True)
    add_footer(s, 5)

    # 6 results
    s = prs.slides.add_slide(blank)
    add_title(s, "复现实验：le450 基准与随机平面图")
    add_text(s, "所有最终结果都用冲突边数校验；合法标准统一为 conflict = 0", 0.8, 1.25, 8.2, 0.3, size=15, color=GRAY)
    add_three_line_table(s, ["数据", "顶点", "边", "颜色", "耗时/s"], [["le450_5a", "450", "5714", "5", "1.104681"], ["le450_15b", "450", "8169", "15", "2.501488"], ["le450_25a", "450", "8260", "25", "0.007555"]], 0.8, 1.8, 5.4, [1.75, 0.75, 0.85, 0.75, 1.3], row_h=0.34, font_size=14)
    add_text(s, "三组均成功，冲突边数均为 0。", 0.8, 3.5, 4.7, 0.3, size=14, color=GRAY)
    add_three_line_table(s, ["n", "边数", "耗时/s", "冲突"], [["20", "54", "0.000280", "0"], ["80", "234", "0.001980", "0"], ["160", "474", "0.005360", "0"], ["300", "894", "0.020500", "0"]], 7.2, 1.8, 4.1, [0.75, 0.85, 1.6, 0.8], row_h=0.34, font_size=14)
    add_text(s, "随机平面图全部 4 色成功。", 7.2, 3.82, 4.1, 0.3, size=14, color=GRAY)
    add_text(s, "颜色数越宽松，局部搜索越容易收敛；边数更多、颜色更少时约束更强。", 0.8, 5.25, 11.5, 0.65, size=23, color=BLUE, bold=True)
    add_footer(s, 6)

    # 7 conclusion
    s = prs.slides.add_slide(blank)
    add_title(s, "5 分钟汇报结论：我做了什么，优化在哪里")
    add_text(s, "结束页只保留可答辩的三点", 0.8, 1.25, 5.0, 0.3, size=15, color=GRAY)
    add_multiline(s, [
        "模型设计：统一把手工地图、DIMACS .col 和随机平面图表示为无向图；用冲突边数 = 0 作为唯一合法性判据。",
        "算法优化：从暴力回溯的固定顺序，改为 DSATUR 的饱和度优先和前向检查；大图使用 TabuCol 避免精确搜索爆炸。",
        "复现结果：小图、随机平面图、三个 le450 数据全部得到 0 冲突解；DSATUR 降低搜索赋值次数，TabuCol 在 450 点图上快速收敛。"
    ], 1.05, 2.0, 11.2, 2.7, size=17, numbered=True)
    add_text(s, "参考：Brélaz, CACM 1979; Robertson–Sanders–Seymour–Thomas, JCTB 1997; Leighton, J. Res. NBS 1979.", 0.8, 5.85, 11.4, 0.5, size=13, color=GRAY)
    add_footer(s, 7)

    prs.save(OUT)


if __name__ == "__main__":
    build()
    print(OUT)
